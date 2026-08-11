# -*- coding: utf-8 -*-
"""CASPAA Strategy & Roadmap deck - 8 weeks, two phases.

Structure follows the Academy of Product Management template (Vision, Strategy
via DIBB, Goals, Roadmap Themes, theme details, timeline).

Ordering is by what a school can actually run on, not by which price tier a
feature sits in. Phase 1 (weeks 1-4) is the set a school needs to get through a
full term without paper. Phase 2 (weeks 5-8) is the parent, teacher and student
apps, plus the features schools ask for next.

Body copy never drops below 12pt. Sections carrying more than that allows are
split across slides rather than shrunk.
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

OUT = r"C:\Users\USER\Desktop\CASPAA\CASPAA_Strategy_Roadmap.pptx"
NAIRA = u"\u20a6"

# ---------------------------------------------------------------- palette
NAVY = RGBColor(0x10, 0x27, 0x30)
TEAL = RGBColor(0x0A, 0x6B, 0x50)
ACC = RGBColor(0x0E, 0x9E, 0x76)
GREY = RGBColor(0x4F, 0x61, 0x67)
BODY = RGBColor(0x28, 0x38, 0x3E)
HOLD = RGBColor(0x8E, 0x4E, 0x5C)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
SOFT = RGBColor(0xED, 0xF3, 0xF1)
PANEL = RGBColor(0xEC, 0xF1, 0xF3)
LINE = RGBColor(0xC7, 0xD2, 0xD6)
GOLD = RGBColor(0xB8, 0x8A, 0x1E)

P1 = RGBColor(0x0A, 0x6B, 0x50)
P2 = RGBColor(0x3D, 0x6E, 0xA8)

T_COLORS = [
    RGBColor(0x10, 0x33, 0x3F),   # Foundations
    RGBColor(0x0E, 0x9E, 0x76),   # The school day
    RGBColor(0xC2, 0x62, 0x2B),   # Fees and payments
    RGBColor(0x0A, 0x6B, 0x50),   # Reaching parents
    RGBColor(0x3D, 0x6E, 0xA8),   # The apps
    RGBColor(0x7A, 0x5E, 0xA8),   # Teaching and learning
    RGBColor(0xB8, 0x8A, 0x1E),   # Staff, money and scale
]

FONT = "Segoe UI"
W, H = 13.333, 7.5
M = 0.62

prs = Presentation()
prs.slide_width = Inches(W)
prs.slide_height = Inches(H)
BLANK = prs.slide_layouts[6]


# ---------------------------------------------------------------- helpers
def slide():
    return prs.slides.add_slide(BLANK)


def textbox(s, x, y, w, h):
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    return tf


def para(tf, text, *, size=14, bold=False, italic=False, color=BODY,
         space_before=0, space_after=5, first=False, align=None, line=1.22):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.space_before = Pt(space_before)
    p.space_after = Pt(space_after)
    p.line_spacing = line
    if align is not None:
        p.alignment = align
    r = p.add_run()
    r.text = text
    f = r.font
    f.name = FONT
    f.size = Pt(size)
    f.bold = bold
    f.italic = italic
    f.color.rgb = color
    return p


def rich(p, text, *, size=14, bold=False, italic=False, color=BODY):
    r = p.add_run()
    r.text = text
    f = r.font
    f.name = FONT
    f.size = Pt(size)
    f.bold = bold
    f.italic = italic
    f.color.rgb = color
    return r


def brand(s):
    tf = textbox(s, W - M - 3.4, 0.32, 3.4, 0.5)
    para(tf, "CASPAA", size=18, bold=True, color=NAVY,
         align=PP_ALIGN.RIGHT, first=True, space_after=0)
    tf2 = textbox(s, W - M - 3.4, 0.68, 3.4, 0.3)
    para(tf2, "SCHOOL OPERATING SYSTEM", size=11, bold=True, color=TEAL,
         align=PP_ALIGN.RIGHT, first=True)


def heading(s, title, sub=None, eyebrow=None, *, title_size=36, sub_w=10.4,
            sub_gap=0.66, eyebrow_color=TEAL):
    brand(s)
    y = 0.5
    if eyebrow:
        tf = textbox(s, M, y, 8.6, 0.3)
        para(tf, eyebrow.upper(), size=11.5, bold=True, color=eyebrow_color, first=True)
        y += 0.34
    tf = textbox(s, M, y, 9.3, 0.8)
    para(tf, title, size=title_size, color=NAVY, first=True, space_after=0)
    y += title_size / 55.0 + 0.24
    if sub:
        tf = textbox(s, M, y, sub_w, 0.8)
        para(tf, sub, size=15, italic=True, color=GREY, first=True, line=1.25)
        y += sub_gap
    return y + 0.22


def box(s, x, y, w, h, fill=PANEL, line_color=None):
    shape = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y),
                               Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if line_color is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(0.75)
    shape.shadow.inherit = False
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Inches(0.16)
    tf.margin_top = tf.margin_bottom = Inches(0.12)
    return shape


def bar(s, x, y, w, h, fill, label, *, size=11.5):
    shape = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y),
                               Inches(w), Inches(h))
    shape.adjustments[0] = 0.24
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.fill.background()
    shape.shadow.inherit = False
    tf = shape.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = tf.margin_right = Inches(0.08)
    tf.margin_top = tf.margin_bottom = 0
    para(tf, label, size=size, bold=True, color=WHITE, first=True,
         space_after=0, align=PP_ALIGN.CENTER)
    return shape


def footer(s, text):
    tf = textbox(s, M, H - 0.55, W - 2 * M - 3.2, 0.34)
    para(tf, text, size=10.5, italic=True, color=GREY, first=True)


def pager(s, text):
    tf = textbox(s, W - M - 3.2, H - 0.55, 3.2, 0.34)
    para(tf, text, size=10.5, bold=True, color=GREY, first=True,
         align=PP_ALIGN.RIGHT)


def style_table(table, *, header_size=13, body_size=12.5):
    for cell in table.rows[0].cells:
        cell.fill.solid()
        cell.fill.fore_color.rgb = NAVY
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        cell.margin_left = cell.margin_right = Inches(0.12)
        for p in cell.text_frame.paragraphs:
            for r in p.runs:
                r.font.name = FONT
                r.font.size = Pt(header_size)
                r.font.bold = True
                r.font.color.rgb = WHITE
    for ri, row in enumerate(list(table.rows)[1:], 1):
        for cell in row.cells:
            cell.fill.solid()
            cell.fill.fore_color.rgb = WHITE if ri % 2 else SOFT
            cell.vertical_anchor = MSO_ANCHOR.TOP
            cell.margin_left = cell.margin_right = Inches(0.12)
            cell.margin_top = cell.margin_bottom = Inches(0.09)
            for p in cell.text_frame.paragraphs:
                p.line_spacing = 1.15
                p.space_after = Pt(3)
                for r in p.runs:
                    r.font.name = FONT
                    if r.font.size is None:
                        r.font.size = Pt(body_size)
                    if r.font.color.type is None:
                        r.font.color.rgb = BODY


def chunk(seq, n):
    return [seq[i:i + n] for i in range(0, len(seq), n)]


# ================================================================ content
WEEKS = [
    ("Wk 1", "3 Aug"), ("Wk 2", "10 Aug"), ("Wk 3", "17 Aug"), ("Wk 4", "24 Aug"),
    ("Wk 5", "31 Aug"), ("Wk 6", "7 Sep"), ("Wk 7", "14 Sep"), ("Wk 8", "21 Sep"),
]
WEEK_END = ["7 Aug", "14 Aug", "21 Aug", "28 Aug",
            "4 Sep", "11 Sep", "18 Sep", "25 Sep"]

THEMES = [
    {
        "name": "Foundations",
        "phase": 1,
        "modules": "Backend & API  ·  Login & RBAC  ·  Tenancy  ·  Offline",
        "desc": "One backend, one login, hard separation between schools, and an app that "
                "keeps working when the network drops. Invisible to anyone using it, and "
                "load-bearing for every theme after it, including the apps.",
        "hypothesis": "Finish the spine inside two weeks and three squads can run in "
                      "parallel for the remaining six without tripping over each other.",
        "validation": "Seven roles sign in to the seed school and each sees only their "
                      "own data. A teacher marks a register with the laptop in flight "
                      "mode and it syncs when the network returns.",
        "actions": [
            {"name": "Production backend and API", "short": "Backend & API",
             "why": "A school's data belongs on a server the apps and the web can both read.",
             "milestones": ["Data model and schema built from the PRD",
                            "REST API with a validation layer, designed for the phone apps to reuse",
                            "A pilot school's records loaded and dogfooded internally"],
             "start": 1, "end": 1, "bar": "Backend & API"},
            {"name": "Unified login and RBAC", "short": "Login & RBAC",
             "why": "Seven roles, one front door. Get it wrong and every screen inherits the mistake.",
             "milestones": ["Role-aware routing from a single login",
                            "Permissions enforced server-side, not in the interface",
                            "Token-based sessions the apps can use later"],
             "start": 1, "end": 1, "bar": "Login & RBAC"},
            {"name": "Tenant isolation", "short": "Tenancy",
             "why": "Complete separation between schools is the assurance every proprietor asks for first.",
             "milestones": ["School scoping enforced on every read and write",
                            "Automated isolation tests in the build"],
             "start": 1, "end": 2, "bar": "Tenant isolation"},
            {"name": "Offline-first behaviour", "short": "Offline",
             "why": "Registers get marked in classrooms with no signal, so the app has to work without one.",
             "milestones": ["Installable app, offline reads",
                            "Register and results queue, then sync on reconnect"],
             "start": 2, "end": 2, "bar": "Offline"},
        ],
    },
    {
        "name": "The school day",
        "phase": 1,
        "modules": "Students  ·  Admissions  ·  Timetable  ·  Register  ·  Results",
        "desc": "Enrol a child, put them in a class, take the register, enter the marks, "
                "approve them, print the report card. This loop is the school's operating "
                "week, and it is what Phase 1 exists to deliver.",
        "hypothesis": "Once a teacher can finish a whole teaching day in the app, the "
                      "paper copy stops being kept. That is when the data becomes good "
                      "enough to bill and report from.",
        "validation": "One class arm runs a full week end to end. Register daily, marks "
                      "entered and approved, a report card printed on the school's own "
                      "template with nobody falling back to a notebook.",
        "actions": [
            {"name": "Students and admissions", "short": "Students & admissions",
             "why": "Every other module reads the student record, so it is the first thing built.",
             "milestones": ["Enrolment, bulk CSV import, online applications",
                            "Promotion between classes and sessions"],
             "start": 1, "end": 2, "bar": "Students & admissions"},
            {"name": "Sessions, classes and timetable", "short": "Classes & timetable",
             "why": "Terms and class arms are what the register, the results and the fees all hang off.",
             "milestones": ["Sessions, terms and class arms",
                            "Timetable with clash detection",
                            "Subject and teacher assignment"],
             "start": 2, "end": 2, "bar": "Classes & timetable"},
            {"name": "Attendance register", "short": "Attendance",
             "why": "Marked 190 days a year, which makes it the most-used screen in the product.",
             "milestones": ["Daily class register that works offline",
                            "Clock-in and clock-out student register",
                            "Automatic parent alert on an unexplained absence"],
             "start": 2, "end": 3, "bar": "Attendance register"},
            {"name": "Results and report cards", "short": "Results",
             "why": "Termly, but it is the document a school is judged on by every parent at once.",
             "milestones": ["Teacher entry, admin approval gate, broadsheet",
                            "Printable report cards carrying the school's branding"],
             "start": 3, "end": 4, "bar": "Results & report cards"},
        ],
    },
    {
        "name": "Fees and payments",
        "phase": 1,
        "modules": "Fee structure  ·  Invoices  ·  Ledger  ·  Payments  ·  Reconciliation",
        "desc": "The reason a proprietor signs the contract. Bill a term, show every "
                "parent an itemised statement, take the money online, and have it "
                "reconcile itself. Our own subscription billing lands here too, priced per "
                "student per term off the live roll.",
        "hypothesis": "Parents pay sooner when they can see exactly what they owe and a "
                      "receipt arrives the moment the money does. Schools renew when the "
                      "collection rate visibly moves.",
        "validation": "The pilot school bills a real term on CASPAA. Payments match "
                      "automatically above 98%, statements and receipts print correctly, "
                      "and the owing list agrees with the bank without a spreadsheet.",
        "actions": [
            {"name": "Fee structure, invoices and the student ledger", "short": "Fees & ledger",
             "why": "The itemised statement of account is the one feature parents ask for by name.",
             "milestones": ["Fee structure and invoice generation per term",
                            "Itemised statement with the discount breakdown",
                            "Owing and part-paid filters, advance credit"],
             "start": 2, "end": 3, "bar": "Fees, invoices & ledger"},
            {"name": "Online payments and receipts", "short": "Payments",
             "why": "Card, transfer and USSD with an instant receipt is what turns billed into collected.",
             "milestones": ["Paystack card, bank transfer and USSD",
                            "Webhooks and real-time receipts"],
             "start": 3, "end": 4, "bar": "Payments & receipts"},
            {"name": "Reconciliation and financial reports", "short": "Reconciliation",
             "why": "Owners decide on billed against collected. Bursars judge us on the match rate.",
             "milestones": ["Automatic bank reconciliation above 98%",
                            "Billed, collected and outstanding by class and term"],
             "start": 4, "end": 4, "bar": "Recon & reports"},
            {"name": "Our own subscription billing", "short": "CASPAA billing",
             "why": "We charge per student per term, so the invoice has to come off the live roll rather than a fixed figure.",
             "milestones": ["One pricing source shared by the website and the app",
                            "Term invoice computed from live headcount, setup fee as a one-off line"],
             "start": 4, "end": 4, "bar": "Subscription billing"},
        ],
    },
    {
        "name": "Reaching parents",
        "phase": 1,
        "modules": "Parent portal  ·  Announcements  ·  Consent  ·  WhatsApp & email",
        "desc": "A school that cannot reach its parents has bought a filing cabinet. The "
                "web portal and the messaging channels come first; the phone app in Phase "
                "2 sits on top of the same work.",
        "hypothesis": "Parents who receive absence alerts and receipts on WhatsApp open "
                      "the portal without being chased, which is what makes fee reminders "
                      "work later on.",
        "validation": "Every parent at the pilot school receives at least one live message "
                      "in week 4, and consent for one real school trip is collected "
                      "entirely in the app.",
        "actions": [
            {"name": "Parent portal and child dashboards", "short": "Parent portal",
             "why": "One login covering several children is the difference between a parent using it and not.",
             "milestones": ["Per-child view of fees, results and attendance",
                            "Statement and receipt download"],
             "start": 3, "end": 3, "bar": "Parent portal"},
            {"name": "Announcements, calendar and digital consent", "short": "Announcements & consent",
             "why": "Consent slips are the paper the front office loses most often.",
             "milestones": ["Announcements and notice board",
                            "School calendar and events",
                            "Digital consent with a live response list"],
             "start": 3, "end": 4, "bar": "Announcements, calendar & consent"},
            {"name": "WhatsApp and email delivery", "short": "WhatsApp & email",
             "why": "In-app notifications alone do not reach Nigerian parents. WhatsApp does.",
             "milestones": ["Provider integration and delivery receipts",
                            "Absence, fee and receipt templates"],
             "start": 4, "end": 4, "bar": "WhatsApp & email"},
        ],
    },
    {
        "name": "The apps",
        "phase": 2,
        "modules": "Parent  ·  Teacher  ·  Student  ·  Store release",
        "desc": "Phones for the three roles that are not sitting at a desk. Parents, "
                "teachers and students. Admin, finance and the operations portal stay on "
                "the web, where the work is genuinely desk work.",
        "hypothesis": "Parents will not keep a browser tab for the school, but they will "
                      "keep an icon. The same is true of students. For teachers it is "
                      "about marking the register where the class actually is.",
        "validation": "Installs on real devices at the pilot school. A teacher marks a "
                      "register on a phone in a classroom with no signal and it syncs "
                      "later. A parent pays fees from the app without opening a browser.",
        "actions": [
            {"name": "App shell, auth and sync", "short": "App shell & sync",
             "why": "One shell serving three roles keeps the apps consistent and the release predictable.",
             "milestones": ["Shared shell with role-aware navigation",
                            "Token auth against the Phase 1 API",
                            "Offline store and background sync"],
             "start": 5, "end": 6, "bar": "Shell, auth & sync"},
            {"name": "Parent app", "short": "Parent app",
             "why": "The highest-volume audience, and the one that pays.",
             "milestones": ["Fees, statement and in-app payment",
                            "Results, attendance and absence alerts",
                            "Consent and announcements with push"],
             "start": 6, "end": 7, "bar": "Parent app"},
            {"name": "Teacher app", "short": "Teacher app",
             "why": "The register belongs in a pocket, not on a desktop in the staff room.",
             "milestones": ["Register marking, offline first",
                            "Results entry and the diary"],
             "start": 6, "end": 7, "bar": "Teacher app"},
            {"name": "Student app and store release", "short": "Student app & release",
             "why": "The app counts once a parent can install it from the store they already use.",
             "milestones": ["Timetable, assignments, results",
                            "Play Store and App Store submission, with the review time budgeted"],
             "start": 7, "end": 8, "bar": "Student app & store release"},
        ],
    },
    {
        "name": "Teaching and learning",
        "phase": 2,
        "modules": "Assessments  ·  CBT  ·  Lesson plans  ·  Behaviour  ·  Alumni",
        "desc": "What a school reaches for once the operating week is running smoothly. "
                "Frequently requested, and well placed to follow the essentials rather "
                "than compete with them.",
        "hypothesis": "Assessments are what pull students into the product daily rather "
                      "than at report-card time.",
        "validation": "One CBT sitting runs for a real class with real marks, and the "
                      "results land in the same broadsheet as everything else.",
        "actions": [
            {"name": "Assignments and CBT", "short": "Assignments & CBT",
             "why": "Homework and online tests are the reason a student opens the app between terms.",
             "milestones": ["Assignments with inline marking and resubmission",
                            "CBT engine and formative tests"],
             "start": 5, "end": 6, "bar": "Assignments & CBT"},
            {"name": "Lesson plans and materials", "short": "Lessons & materials",
             "why": "Heads of department want visibility of what is being taught, not just what was marked.",
             "milestones": ["Lesson plans against the scheme of work",
                            "Notes and materials for the student learning area"],
             "start": 6, "end": 7, "bar": "Lesson plans & materials"},
            {"name": "Behaviour, discipline and alumni", "short": "Behaviour & alumni",
             "why": "Completes the student record: conduct alongside marks, and what happens after they leave.",
             "milestones": ["Incidents, merits and sanctions",
                            "Alumni register and leaving certificates"],
             "start": 7, "end": 7, "bar": "Behaviour & alumni"},
        ],
    },
    {
        "name": "Staff, money and scale",
        "phase": 2,
        "modules": "Payroll & HR  ·  Accounting  ·  Transport  ·  Inventory  ·  Groups",
        "desc": "Back office and growth. Payroll and the accounting package are what stop "
                "a school keeping a second set of books. Multi-branch closes the plan, "
                "sized for the group owners who arrive once the single campus is proven.",
        "hypothesis": "Payroll is the strongest reason an existing school expands its "
                      "usage rather than churning. Multi-branch raises the price but "
                      "rarely wins the first deal.",
        "validation": "One payroll cycle runs end to end at the pilot school, and its "
                      "accountant accepts the trial balance without re-keying it. A group "
                      "owner runs two campuses from one login.",
        "actions": [
            {"name": "Payroll, HR and salary advance", "short": "Payroll & HR",
             "why": "The clearest reason a school widens what it uses instead of drifting away.",
             "milestones": ["Payroll run and payslip distribution",
                            "Leave, appraisals, substitute cover",
                            "Salary advance against payroll"],
             "start": 5, "end": 6, "bar": "Payroll, HR & advances"},
            {"name": "Accounting package", "short": "Accounting",
             "why": "P&L, trial balance and balance sheet are what stop the second set of books.",
             "milestones": ["Chart of accounts and posting rules",
                            "P&L, trial balance, balance sheet",
                            "An export the external accountant will accept"],
             "start": 6, "end": 7, "bar": "Accounting package"},
            {"name": "Transport and authorised pickup", "short": "Transport",
             "why": "Safety carries more weight with parents than anything on the academic side.",
             "milestones": ["Routes, assignments, live bus status",
                            "Authorised pickup approval"],
             "start": 6, "end": 7, "bar": "Transport & pickup"},
            {"name": "Inventory, house points and surveys", "short": "Operations",
             "why": "What the bursar and the sports master ask for by the second term.",
             "milestones": ["Inventory with low-stock alerts",
                            "House points, leaderboards, inter-house competitions",
                            "Surveys and feedback forms"],
             "start": 7, "end": 8, "bar": "Inventory, houses & surveys"},
            {"name": "Multi-branch groups", "short": "Multi-branch",
             "why": "One login across campuses, and the feature a group owner signs for.",
             "milestones": ["Group overview and branch switcher",
                            "Add-branch on-ramp",
                            "Consolidated finance and enrolment roll-up"],
             "start": 8, "end": 8, "bar": "Multi-branch groups"},
        ],
    },
]

PH1 = [t for t in THEMES if t["phase"] == 1]
PH2 = [t for t in THEMES if t["phase"] == 2]


# ---------------------------------------------------------------- 1. title
s = slide()
brand(s)
tf = textbox(s, M, 2.5, 12.0, 1.3)
para(tf, "Strategy & Roadmap", size=58, color=NAVY, first=True, space_after=0)
tf = textbox(s, M, 3.95, 12.0, 1.1)
para(tf, "CASPAA School Operating System", size=24, bold=True, color=TEAL,
     first=True, space_after=4)
para(tf, "Built new in eight weeks. What a school runs on first, then the apps.",
     size=22, color=GREY)
tf = textbox(s, M, 5.45, 12.0, 0.5)
para(tf, "3 August to 25 September 2026   ·   Owner: Product   ·   Draft for stakeholder review",
     size=14, italic=True, color=GREY, first=True)

# ---------------------------------------------------------------- 2. vision
s = slide()
y = heading(s, "Vision", "Where this ends up if the strategy works.")
b = box(s, M, y, W - 2 * M, 1.3, fill=SOFT)
tf = b.text_frame
para(tf, "A school runs its whole day on CASPAA and nobody in the building keeps a "
         "parallel paper record. That is the test.",
     size=22, bold=True, color=NAVY, first=True, line=1.18)
y += 1.58

cols = [
    ("The school", "Enrolment, the register, results and fees in one place. The bursar "
                   "stops reconciling by hand. The head sees the collection rate without "
                   "having to ask for it."),
    ("The parent", "An icon on the phone that answers both questions: what do I owe, and "
                   "how is my child doing. Statement, receipt, report card, and a way to "
                   "reach the teacher."),
    ("CASPAA", "A vertical operating system rather than a module vendor. Priced per "
               "student, so revenue grows with the schools we serve and not only with "
               "the number we sign."),
]
cw = (W - 2 * M - 0.5) / 3
for i, (t, d) in enumerate(cols):
    x = M + i * (cw + 0.25)
    bx = box(s, x, y, cw, 2.3, fill=WHITE, line_color=LINE)
    tf = bx.text_frame
    para(tf, t, size=17, bold=True, color=TEAL, first=True, space_after=8)
    para(tf, d, size=14, color=BODY, line=1.28)
y += 2.55
tf = textbox(s, M, y, W - 2 * M, 0.5)
p = para(tf, "Eight weeks gets us the first two steps. ", size=14, bold=True,
         color=NAVY, first=True)
rich(p, "A school running a full term on the system, and apps in the hands of the "
        "parents, teachers and students who are never at a desk.", size=14, color=BODY)

# ------------------------------------------------------- 3-4. strategy (DIBB)
DIBB = [
    ("Data", "What we know", [
        "Schools run the day on paper registers, spreadsheets and WhatsApp. Who has paid "
        "and who owes is a question with no fast answer.",
        "Attendance, results and report cards are computed by hand, and the errors "
        "surface a term later in front of parents.",
        "The incumbent is spreadsheets and WhatsApp, not a rival product. Free, "
        "fragmented, no audit trail.",
        "Connections drop in most of the buildings we sell into. Cloud-only tools stop "
        "at the door.",
        "Tens of thousands of Nigerian private schools are digitising for the first time.",
    ]),
    ("Insights", "What we learned", [
        "Schools are not short of records. They are short of the operating loop that "
        "produces those records without hand-work. Managing records and running a school "
        "are different products.",
        "Fees is the problem a proprietor will pay to solve this term. Everything else we "
        "sell is bought on the back of it.",
        "Parents and students live on a phone; office staff live at a desk. One product "
        "shaped for both surfaces means the phone loses.",
        "Offline is a condition of the building, not a feature of the software.",
    ]),
    ("Beliefs", "What we think wins", [
        "The operating week is the product: enrol, timetable, register, marks, invoice, "
        "receipt, message the parent. Everything else is an addition to it.",
        "Whoever runs a school's operating layer earns the payment flow on top of it. "
        "That order does not reverse.",
        "A school onboards once, and judges us on what it finds in that first week. "
        "Fewer modules working beats more modules partly built.",
        "Trust is built on the fee conversation. An itemised statement a parent can read "
        "does more for retention than any academic feature.",
    ]),
    ("Bets", "What we are backing", [
        "Build new against the specification rather than growing the prototype. The "
        "prototype proved the shape of the product; it was never the foundation.",
        "Four weeks to a product a school can run a term on, then four weeks on the apps "
        "for parents, teachers and students.",
        "One app shell across the three roles, with the interface changing by role.",
        "Put it into one pilot school at week 4 and let that school's term set the order "
        "of Phase 2.",
    ]),
]
for page, pair in enumerate(chunk(DIBB, 2), 1):
    s = slide()
    names = " and ".join(t for t, _, _ in pair)
    y = heading(s, "Strategy",
                "Data, insights, beliefs, bets. This slide covers %s." % names)
    cw = (W - 2 * M - 0.4) / 2
    for i, (t, sub, items) in enumerate(pair):
        idx = [d[0] for d in DIBB].index(t)
        x = M + i * (cw + 0.4)
        bx = box(s, x, y, cw, 4.45, fill=WHITE, line_color=LINE)
        tf = bx.text_frame
        para(tf, t, size=22, bold=True, color=T_COLORS[idx], first=True, space_after=2)
        para(tf, sub, size=12.5, italic=True, color=GREY, space_after=12)
        for it in items:
            p = para(tf, "", space_after=9, line=1.25)
            rich(p, "·   ", size=13.5, bold=True, color=ACC)
            rich(p, it, size=13.5, color=BODY)
    pager(s, "Strategy  %d of 2" % page)

# ------------------------------------------------------ 5. how we prioritised
s = slide()
y = heading(s, "How we chose the order",
            "One question decides Phase 1: can a school get through a full term on it?")
RULES = [
    ("1", "If a school cannot complete a term without it, it is Phase 1.",
     "Enrol, timetable, register, results, report cards, invoices, payments, tell the parents."),
    ("2", "Daily beats termly, and termly beats annual.",
     "The register runs 190 days a year. Report cards run three times. Alumni records, once."),
    ("3", "Anything that blocks money coming in outranks everything except the spine.",
     "Invoices and reconciliation are why the proprietor signed, and how we get paid ourselves."),
    ("4", "Ship what we can finish.",
     "Scope comes out before the date moves. A half-built module in a live school costs more than a missing one."),
    ("5", "The apps ship on a real API.",
     "Phase 1 delivers the API the apps consume, which is what lets app work start on solid ground in week 5."),
]
rowy = y
for n, rule, why in RULES:
    bx = box(s, M, rowy, W - 2 * M, 0.7, fill=WHITE, line_color=LINE)
    tf = bx.text_frame
    p = para(tf, "", first=True, space_after=0, line=1.2)
    rich(p, n + ".   ", size=15, bold=True, color=ACC)
    rich(p, rule + "   ", size=15, bold=True, color=NAVY)
    rich(p, why, size=12.5, color=GREY)
    rowy += 0.78
bx = box(s, M, rowy + 0.04, W - 2 * M, 0.8, fill=SOFT)
tf = bx.text_frame
p = para(tf, "", first=True, line=1.24)
rich(p, "The line we hold. ", size=13.5, bold=True, color=TEAL)
rich(p, "Phase 1 is finished when a school can complete a term on it. Everything that "
        "does not serve that moves to Phase 2, whichever plan it is sold in and however "
        "easy it would have been to add.", size=13.5, color=BODY)

# ---------------------------------------------------------------- 6. goals
s = slide()
y = heading(s, "Goals", "Two gates, four weeks apart. Baseline is today, 29 July 2026.")
GATES = [
    ("Gate 1 - 28 August: a school can run on it", TEAL, [
        "One pilot school completes a full week with no paper fallback",
        "Register taken in-app by 90% of classes, every day",
        "A real term billed, collected and reconciled above 98%",
        "Every parent reached at least once on a live channel",
    ]),
    ("Gate 2 - 25 September: the apps are in hand", P2, [
        "Parent, teacher and student apps live in both stores",
        "Half the pilot school's parents installed and signed in",
        "A register marked on a phone, offline, syncing later",
        "Assessments, payroll and accounting live for the schools that want them",
    ]),
]
cw = (W - 2 * M - 0.4) / 2
for i, (title, colour, items) in enumerate(GATES):
    x = M + i * (cw + 0.4)
    bx = box(s, x, y, cw, 2.4, fill=WHITE, line_color=LINE)
    tf = bx.text_frame
    para(tf, title, size=16, bold=True, color=colour, first=True, space_after=10)
    for it in items:
        p = para(tf, "", space_after=7, line=1.22)
        rich(p, "·  ", size=13.5, bold=True, color=colour)
        rich(p, it, size=13.5, color=BODY)
y += 2.65
tf = textbox(s, M, y, 7.0, 0.32)
para(tf, "THE YEAR-ONE NUMBERS THESE SERVE", size=11.5, bold=True, color=TEAL, first=True)
y += 0.38
KPI = [("60", "schools onboarded"), (NAIRA + "1B+", "payments processed"),
       ("85%", "fee-collection rate"), ("99.9%", "platform uptime")]
kw = (W - 2 * M - 3 * 0.3) / 4
for i, (v, lab) in enumerate(KPI):
    x = M + i * (kw + 0.3)
    bx = box(s, x, y, kw, 1.05, fill=SOFT)
    tf = bx.text_frame
    para(tf, v, size=28, bold=True, color=NAVY, first=True, space_after=2)
    para(tf, lab, size=13, color=GREY, space_after=0)

# ------------------------------------------------- roadmap themes
THEME_PAGES = [(PH1[0], PH1[1]), (PH1[2], PH1[3]), (PH2[0], PH2[1]), (PH2[2],)]
for page, pair in enumerate(THEME_PAGES, 1):
    s = slide()
    y = heading(s, "Roadmap Themes",
                "Phase 1, weeks 1 to 4, is what a school runs on. Phase 2, weeks 5 to 8, "
                "is the apps and the next layer.", title_size=32, sub_gap=0.46)
    cw = (W - 2 * M - 0.4) / 2
    for j, th in enumerate(pair):
        i = THEMES.index(th)
        x = M + j * (cw + 0.4)
        tf = textbox(s, x, y, cw, 0.4)
        p = para(tf, th["name"], size=20, bold=True, color=T_COLORS[i], first=True,
                 space_after=0)
        rich(p, "   Phase %d" % th["phase"], size=13, bold=True,
             color=P1 if th["phase"] == 1 else P2)
        tf = textbox(s, x, y + 0.38, cw, 0.3)
        para(tf, th["modules"], size=11.5, bold=True, color=GREY, first=True)
        bx = box(s, x, y + 0.7, cw, 4.0, fill=PANEL)
        tf = bx.text_frame
        para(tf, th["desc"], size=13, color=BODY, first=True, space_after=10, line=1.25)
        p = para(tf, "", space_after=8, line=1.25)
        rich(p, "Hypothesis. ", size=13, bold=True, color=NAVY)
        rich(p, th["hypothesis"], size=13, italic=True, color=GREY)
        p = para(tf, "", space_after=8, line=1.25)
        rich(p, "How we validate it. ", size=13, bold=True, color=NAVY)
        rich(p, th["validation"], size=13, italic=True, color=GREY)
        p = para(tf, "", space_after=3, line=1.25)
        rich(p, "Actions: ", size=13, bold=True, color=NAVY)
        rich(p, ", ".join(a["short"] for a in th["actions"]), size=13, color=BODY)
    pager(s, "Roadmap Themes  %d of %d" % (page, len(THEME_PAGES)))

# --------------------------------------- theme details (3 actions per slide)
for i, th in enumerate(THEMES):
    pages = chunk(th["actions"], 3)
    for pi, actions in enumerate(pages, 1):
        s = slide()
        y = heading(s, th["name"], eyebrow="Theme details  ·  Phase %d" % th["phase"],
                    title_size=34, eyebrow_color=P1 if th["phase"] == 1 else P2)
        tf = textbox(s, M, y - 0.06, W - 2 * M, 0.36)
        para(tf, th["modules"], size=13, bold=True, color=T_COLORS[i], first=True)
        y += 0.44

        tbl = s.shapes.add_table(len(actions) + 1, 4, Inches(M), Inches(y),
                                 Inches(W - 2 * M), Inches(3.1)).table
        tbl.columns[0].width = Inches(2.9)
        tbl.columns[1].width = Inches(3.2)
        tbl.columns[2].width = Inches(4.3)
        tbl.columns[3].width = Inches(1.69)
        for c, h in enumerate(["Action", "Why it matters", "Milestones (epics)", "ETA"]):
            tbl.cell(0, c).text = h
        for r, a in enumerate(actions, 1):
            tbl.cell(r, 0).text = a["name"]
            tbl.cell(r, 1).text = a["why"]
            mc = tbl.cell(r, 2).text_frame
            for j, ms in enumerate(a["milestones"]):
                p = mc.paragraphs[0] if j == 0 else mc.add_paragraph()
                p.line_spacing = 1.15
                p.space_after = Pt(3)
                rich(p, "·  ", size=12.5, bold=True, color=T_COLORS[i])
                rich(p, ms, size=12.5, color=BODY)
            span = ("Week %d" % a["start"] if a["start"] == a["end"]
                    else "Weeks %d-%d" % (a["start"], a["end"]))
            ec = tbl.cell(r, 3).text_frame
            rich(ec.paragraphs[0], span, size=13, bold=True, color=T_COLORS[i])
            p2 = ec.add_paragraph()
            p2.space_after = Pt(0)
            rich(p2, "by " + WEEK_END[a["end"] - 1], size=11.5, color=GREY)
        style_table(tbl, header_size=13, body_size=12.5)
        for r in range(1, len(actions) + 1):
            for run_ in tbl.cell(r, 0).text_frame.paragraphs[0].runs:
                run_.font.bold = True
                run_.font.size = Pt(13.5)
                run_.font.color.rgb = NAVY
            for run_ in tbl.cell(r, 1).text_frame.paragraphs[0].runs:
                run_.font.italic = True
                run_.font.color.rgb = GREY
        footer(s, "Done means tenant-scoped, role-correct, all states handled, responsive, "
                  "audited, offline-tolerant, acceptance criteria signed off.")
        if len(pages) > 1:
            pager(s, "%s  %d of %d" % (th["name"], pi, len(pages)))


# ------------------------------------------------ timeline helpers
def week_header(s, hy, label_w, band=True):
    hx = M
    track_x = hx + label_w
    track_w = W - M - track_x
    colw = track_w / 8
    if band:
        b1 = box(s, track_x, hy - 0.42, colw * 4 - 0.04, 0.36, fill=P1)
        tf = b1.text_frame
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        para(tf, "PHASE 1  ·  A SCHOOL CAN RUN ON IT", size=11.5, bold=True, color=WHITE,
             first=True, space_after=0, align=PP_ALIGN.CENTER)
        b2 = box(s, track_x + colw * 4, hy - 0.42, colw * 4 - 0.04, 0.36, fill=P2)
        tf = b2.text_frame
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        para(tf, "PHASE 2  ·  THE APPS, AND WHAT COMES NEXT", size=11.5, bold=True,
             color=WHITE, first=True, space_after=0, align=PP_ALIGN.CENTER)
    hdr = box(s, hx, hy, label_w, 0.44, fill=RGBColor(0x1A, 0x1B, 0x3A))
    tf = hdr.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    para(tf, "3 Aug - 25 Sep 2026", size=12.5, bold=True, color=WHITE, first=True,
         space_after=0)
    for k, (wk, dt) in enumerate(WEEKS):
        c = box(s, track_x + k * colw, hy, colw - 0.04, 0.44, fill=WHITE,
                line_color=LINE)
        tf = c.text_frame
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        para(tf, wk, size=12, bold=True, color=NAVY, first=True, space_after=0,
             align=PP_ALIGN.CENTER, line=1.05)
        para(tf, dt, size=10.5, color=GREY, space_after=0, align=PP_ALIGN.CENTER,
             line=1.05)
    return track_x, colw


# ------------------------------------------------ roadmap overview
s = slide()
brand(s)
tf = textbox(s, M, 0.42, 9.0, 0.6)
para(tf, "The eight weeks at a glance", size=32, color=NAVY, first=True, space_after=0)
label_w = 3.3
track_x, colw = week_header(s, 1.5, label_w)
y = 1.5 + 0.44 + 0.2
for i, th in enumerate(THEMES):
    st = min(a["start"] for a in th["actions"])
    en = max(a["end"] for a in th["actions"])
    tf = textbox(s, M, y + 0.09, label_w - 0.22, 0.4)
    para(tf, th["name"], size=13.5, bold=True, color=BODY, first=True, space_after=0,
         align=PP_ALIGN.RIGHT)
    bx0 = track_x + (st - 1) * colw
    bw = (en - st + 1) * colw - 0.12
    bar(s, bx0 + 0.03, y, bw, 0.42, T_COLORS[i],
        "%d actions  ·  weeks %d-%d" % (len(th["actions"]), st, en), size=11.5)
    y += 0.55
tf = textbox(s, M, y + 0.1, W - 2 * M, 0.5)
p = para(tf, "", first=True, line=1.25)
rich(p, "Week 4 is a gate. ", size=13.5, bold=True, color=NAVY)
rich(p, "Phase 2 opens once a real school has run a real week on Phase 1, so the app work "
        "starts against a product that is already holding up in a classroom.",
     size=13.5, color=BODY)
pager(s, "Roadmap  overview")

# ------------------------------------------- roadmap detail
DETAIL_PAGES = [(PH1[0], PH1[1]), (PH1[2], PH1[3]), (PH2[0], PH2[1]), (PH2[2],)]
for page, pair in enumerate(DETAIL_PAGES, 1):
    s = slide()
    brand(s)
    tf = textbox(s, M, 0.42, 9.6, 0.6)
    para(tf, "Eight-week roadmap", size=32, color=NAVY, first=True, space_after=0)
    label_w = 2.75
    track_x, colw = week_header(s, 1.55, label_w)
    y = 1.55 + 0.44 + 0.18
    row_h, gap = 0.3, 0.05
    for th in pair:
        i = THEMES.index(th)
        n = len(th["actions"])
        block_h = 0.34 + n * (row_h + gap) + 0.06
        box(s, M, y, W - 2 * M, block_h, fill=PANEL)
        tf = textbox(s, M + 0.16, y + 0.06, 6.0, 0.28)
        para(tf, th["name"].upper(), size=11.5, bold=True, color=T_COLORS[i],
             first=True, space_after=0)
        ry = y + 0.36
        for a in th["actions"]:
            tf = textbox(s, M + 0.14, ry + 0.05, label_w - 0.3, row_h)
            para(tf, a["short"], size=11, bold=True, color=BODY, first=True,
                 space_after=0, align=PP_ALIGN.RIGHT)
            bx0 = track_x + (a["start"] - 1) * colw
            bw = (a["end"] - a["start"] + 1) * colw - 0.12
            # The row is only 0.3in tall, so a caption that would wrap is dropped
            # rather than clipped. The name is on the left either way.
            caption = ""
            for candidate in (a["bar"], a["short"]):
                if len(candidate) * 11 * 0.55 / 72.0 <= bw - 0.18:
                    caption = candidate
                    break
            bar(s, bx0 + 0.03, ry, bw, row_h, T_COLORS[i], caption, size=11)
            ry += row_h + gap
        y += block_h + 0.12
    tf = textbox(s, M, y + 0.06, W - 2 * M, 0.4)
    para(tf, "Each bar is the working window for one action. Names sit on the left.",
         size=11.5, italic=True, color=GREY, first=True)
    pager(s, "Roadmap detail  %d of %d" % (page, len(DETAIL_PAGES)))

# ------------------------------------------------------- what we are not building
s = slide()
y = heading(s, "What we are not building",
            "Worth stating plainly, because each of these has been asked for at least once.")
cw = (W - 2 * M - 0.5) / 3
DEF = [
    ("Held by decision", HOLD, [
        "Fee lending: applications, credit scoring, approval, disbursement, repayment",
        "It ships gated off and stays documented so it can resume without a rewrite",
    ]),
    ("Not in these eight weeks", GREY, [
        "An admin or finance app. That work is desk work and stays on the web",
        "AI report comments and academic insights",
        "Advanced analytics and scheduled reports for the operations portal",
    ]),
    ("Later, if the data supports it", GREY, [
        "Bulk disbursement and NIBSS direct settlement",
        "Savings and investment products, which need a licence we do not hold",
        "AI underwriting, which needs a far larger dataset than we will have this year",
    ]),
]
for i, (title, colour, items) in enumerate(DEF):
    x = M + i * (cw + 0.25)
    bx = box(s, x, y, cw, 2.85, fill=WHITE, line_color=LINE)
    tf = bx.text_frame
    para(tf, title, size=16, bold=True, color=colour, first=True, space_after=10)
    for it in items:
        p = para(tf, "", space_after=8, line=1.25)
        rich(p, "·  ", size=13.5, bold=True, color=colour)
        rich(p, it, size=13.5, color=BODY)
y += 3.1
bx = box(s, M, y, W - 2 * M, 0.85, fill=SOFT)
tf = bx.text_frame
p = para(tf, "", first=True, line=1.24)
rich(p, "Phase 2 is reviewed at the week-4 gate. ", size=14, bold=True, color=NAVY)
rich(p, "Its order is set against what the first live schools are reaching for, which we "
        "will know by then and do not know now.", size=14, color=BODY)

# --------------------------------------------- assumptions & risks
s = slide()
y = heading(s, "What has to hold, and what would move the date")
cw = (W - 2 * M - 0.4) / 2
tf = textbox(s, M, y, cw, 0.32)
para(tf, "ASSUMPTIONS", size=11.5, bold=True, color=TEAL, first=True)
tf = textbox(s, M + cw + 0.4, y, cw, 0.32)
para(tf, "RISKS", size=11.5, bold=True, color=HOLD, first=True)
y += 0.4

A = [("Three squads from week 3",
      "The school day, fees and parent-facing work have no hard dependency on each other "
      "once the spine is done."),
     ("One app shell across the three roles",
      "Parent, teacher and student share a shell and differ by role, which is what keeps "
      "the app work inside four weeks."),
     ("The prototype is the reference, not the foundation",
      "PRD v4 and the working prototype specify the behaviour. These eight weeks build "
      "the product new against them, starting from schema and API in week 1.")]
R = [("Week 2 slips",
      "The spine is unfinished, nothing else starts cleanly and the apps lose their API. "
      "Weeks 1 and 2 carry no other scope so the team can absorb an overrun."),
     ("Store review eats week 8",
      "Play and App Store review is outside our control and can take days. Submission is "
      "planned for week 7, not week 8."),
     ("Done gets traded for speed",
      "Half-finished work in a live school costs more than a missing feature, so done is "
      "enforced at the merge.")]
for col, items, colour in ((0, A, TEAL), (1, R, HOLD)):
    x = M + col * (cw + 0.4)
    yy = y
    for t, d in items:
        bx = box(s, x, yy, cw, 1.32, fill=WHITE, line_color=LINE)
        tf = bx.text_frame
        para(tf, t, size=15, bold=True, color=colour, first=True, space_after=6)
        para(tf, d, size=13, color=BODY, line=1.25)
        yy += 1.46
y += 3 * 1.46
bx = box(s, M, y - 0.08, W - 2 * M, 0.8, fill=SOFT)
tf = bx.text_frame
p = para(tf, "", first=True, line=1.24)
rich(p, "Decision needed before week 5: ", size=14, bold=True, color=NAVY)
rich(p, "package the existing offline web app for the stores, or build the apps native. "
        "Packaging is what fits the four weeks in this plan. Native is a larger piece of "
        "work and needs its own schedule.", size=14, color=BODY)

prs.save(OUT)
print("Wrote", OUT, "-", len(prs.slides._sldIdLst), "slides")
